from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, List, Any
import os
from PIL import Image

class PowerPointGenerator:
    """Генератор PowerPoint презентаций из структуры и дизайна"""
    
    def __init__(self, template_path: str = None):
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        # Стандартные макеты
        self.layouts = {
            "title_slide": 0,
            "title_and_content": 1,
            "section_header": 2,
            "two_content": 3,
            "comparison": 4,
            "title_only": 5,
            "blank": 6
        }
    
    def create_presentation(self, presentation_data: Dict[str, Any], 
                           output_path: str) -> str:
        """
        Создает PowerPoint презентацию из данных
        
        Args:
            presentation_data: Данные от координатора
            output_path: Путь для сохранения
        """
        # Титульный слайд
        self._add_title_slide(presentation_data["topic"])
        
        # Основные слайды
        for i, (structure, content, design) in enumerate(zip(
            presentation_data["structure"],
            presentation_data["content"],
            presentation_data["design"]
        )):
            self._add_content_slide(structure, content, design, i)
        
        # Сохраняем
        self.prs.save(output_path)
        return output_path
    
    def _add_title_slide(self, topic: str):
        """Добавляет титульный слайд"""
        slide_layout = self.prs.slide_layouts[self.layouts["title_slide"]]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = topic
        subtitle.text = "Сгенерировано AI на основе ваших документов"
        
        # Настройка стилей
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
    
    def _add_content_slide(self, structure: Dict, content: Dict, 
                           design: Dict, index: int):
        """Добавляет слайд с контентом"""
        # Выбираем макет на основе дизайна
        layout_id = self.layouts.get(
            design.get("layout_id", "title_and_content"),
            self.layouts["title_and_content"]
        )
        
        slide_layout = self.prs.slide_layouts[layout_id]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Заголовок
        title = slide.shapes.title
        if title:
            title.text = content.get("final_title", structure["title"])
            
            # Применяем цвета из дизайна
            if "color_scheme" in design:
                for paragraph in title.text_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor.from_string(
                        design["color_scheme"].get("primary", "#2E4053")
                    )
        
        # Контент
        if layout_id == self.layouts["title_and_content"]:
            content_placeholder = slide.placeholders[1]
            if content_placeholder:
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                # Добавляем bullet points
                for bullet in content.get("bullets", []):
                    p = text_frame.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(24)
                    p.level = 0
                
                # Добавляем вывод если есть
                if content.get("conclusion"):
                    p = text_frame.add_paragraph()
                    p.text = f"\n{content['conclusion']}"
                    p.font.size = Pt(20)
                    p.font.italic = True
        
        # Добавляем изображения из дизайна
        for visual in design.get("selected_visuals", []):
            if visual.get("source") and os.path.exists(visual["source"]):
                try:
                    left = Inches(10) if visual.get("placement") == "right" else Inches(1)
                    top = Inches(2)
                    
                    # Добавляем картинку
                    slide.shapes.add_picture(
                        visual["source"],
                        left, top,
                        height=Inches(4)
                    )
                except:
                    pass